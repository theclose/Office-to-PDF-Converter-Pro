"""
Shared fixtures for integration + stress tests.
Creates real Office files programmatically for testing.
"""

import os
import tempfile
import pytest


# ─── PDF fixtures (PyMuPDF) ───

def create_test_pdf(pages=3, with_image=False):
    """Create a minimal test PDF with PyMuPDF."""
    try:
        import fitz
    except ImportError:
        pytest.skip("PyMuPDF not available")

    doc = fitz.open()
    for i in range(pages):
        page = doc.new_page(width=595, height=842)  # A4
        page.insert_text((72, 100), f"Test page {i+1} — Trang thử nghiệm" * 20, fontsize=12)
        if with_image:
            rect = fitz.Rect(100, 200, 400, 500)
            page.draw_rect(rect, color=(1, 0, 0), fill=(0.5, 0.5, 1))

    fd, path = tempfile.mkstemp(suffix=".pdf")
    os.close(fd)
    doc.save(path)
    doc.close()
    return path


# ─── Word fixtures (python-docx) ───

def create_test_docx(pages=1):
    """Create a minimal Word document."""
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        pytest.skip("python-docx not available")

    doc = Document()
    for i in range(pages):
        doc.add_paragraph(f"Test page {i+1} — Trang thử nghiệm " * 50)
        if i < pages - 1:
            doc.add_page_break()

    fd, path = tempfile.mkstemp(suffix=".docx")
    os.close(fd)
    doc.save(path)
    return path


# ─── Excel fixtures (openpyxl) ───

def create_test_xlsx(sheets=1, rows=100):
    """Create a minimal Excel workbook."""
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl not available")

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    for r in range(1, rows + 1):
        for c in range(1, 6):
            ws.cell(row=r, column=c, value=f"Data R{r}C{c}")

    for s in range(2, sheets + 1):
        ws2 = wb.create_sheet(title=f"Sheet{s}")
        for r in range(1, 20):
            ws2.cell(row=r, column=1, value=f"Sheet{s} Row{r}")

    fd, path = tempfile.mkstemp(suffix=".xlsx")
    os.close(fd)
    wb.save(path)
    return path


# ─── PowerPoint fixtures (python-pptx) ───

def create_test_pptx(slides=3):
    """Create a minimal PowerPoint presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        pytest.skip("python-pptx not available")

    prs = Presentation()
    for i in range(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = f"Content for slide {i+1}\nTest data — Dữ liệu thử nghiệm"

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


# ─── Pytest fixtures ───

@pytest.fixture
def integration_pdf():
    """Fixture: test PDF, auto-cleanup."""
    path = create_test_pdf(pages=3)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture
def integration_pdf_same_path():
    """Fixture: test PDF for same-path operations (no separate output)."""
    path = create_test_pdf(pages=3, with_image=True)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture
def integration_docx():
    """Fixture: test Word doc, auto-cleanup."""
    path = create_test_docx(pages=2)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture
def integration_xlsx():
    """Fixture: test Excel, auto-cleanup."""
    path = create_test_xlsx(sheets=2, rows=50)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture
def integration_pptx():
    """Fixture: test PPT, auto-cleanup."""
    path = create_test_pptx(slides=3)
    yield path
    if os.path.exists(path):
        os.remove(path)


@pytest.fixture
def temp_output():
    """Fixture: temp output path, auto-cleanup."""
    fd, path = tempfile.mkstemp(suffix=".pdf")
    os.close(fd)
    yield path
    if os.path.exists(path):
        os.remove(path)
